from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parents[1]
OUT_PPTX = BASE_DIR / "docs" / "Negotiator_Pactum_15min.pptx"

FONT_DISPLAY = "Aptos Display"
FONT_BODY = "Aptos"
FONT_MONO = "Consolas"


def rgb(value: str) -> RGBColor:
    value = value.replace("#", "")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


NAVY = rgb("18253C")
INK = rgb("1F2937")
MUTED = rgb("6B7280")
STONE = rgb("9CA3AF")
PAPER = rgb("F6F3EE")
WHITE = rgb("FFFFFF")
SAND = rgb("DED6CA")
GRID = rgb("DDD6CC")
MIST = rgb("EEF2F7")
BLUE = rgb("2563EB")
BLUE_SOFT = rgb("E7F0FF")
TEAL = rgb("0F766E")
TEAL_SOFT = rgb("DDF4F1")
GREEN = rgb("2F855A")
GREEN_SOFT = rgb("DDF4E4")
AMBER = rgb("C47A0A")
AMBER_SOFT = rgb("F9E7C6")
ROSE = rgb("C24163")
ROSE_SOFT = rgb("F9DFE7")
NAVY_SOFT = rgb("E8EDF6")


STRATEGY_ROWS = [
    ("Boulware", "0.7583", "5.00", GREEN),
    ("Baseline", "0.7340", "3.00", BLUE),
    ("Tit for Tat", "0.7340", "3.00", TEAL),
    ("Meso", "0.6922", "5.33", AMBER),
    ("Conceder", "0.5938", "2.25", ROSE),
]

STRATEGY_POINTS = [
    {"name": "Conceder", "round": 2.25, "utility": 0.5938, "color": ROSE},
    {"name": "Baseline", "round": 3.00, "utility": 0.7340, "color": BLUE},
    {"name": "Tit for Tat", "round": 3.00, "utility": 0.7340, "color": TEAL},
    {"name": "Boulware", "round": 5.00, "utility": 0.7583, "color": GREEN},
    {"name": "Meso", "round": 5.33, "utility": 0.6922, "color": AMBER},
]

SCENARIO_MATRIX = [
    ("Margin hardliner", ["R8", "R8", "R8", "R8", "R8"]),
    ("Price floor tradeoff", ["R2", "R8", "R2", "R2", "R2"]),
    ("Payment ceiling", ["R2", "R8", "R2", "A2", "R2"]),
    ("Delivery floor", ["R8", "A8", "R8", "A5", "R8"]),
    ("Late closer", ["R2", "R8", "R2", "R2", "R2"]),
    ("Near settlement", ["A3", "A4", "A5", "A1", "A3"]),
    ("Deadline settlement", ["A3", "A4", "A5", "A1", "A3"]),
]

PRICE_RESULTS = [
    ("Boulware", 106.50, "R7", "R5", GREEN),
    ("Baseline", 111.00, "R7", "R6", BLUE),
    ("Tit for Tat", 113.33, "R7", "R6", TEAL),
    ("Meso", 113.40, "R7", "R5", AMBER),
    ("Conceder", 120.00, "R7", "R5", ROSE),
]

SCRIPT_STEPS = [
    ("T1", "P120 D30 C24", BLUE_SOFT, BLUE),
    ("T2", "P120 D21 C24", TEAL_SOFT, TEAL),
    ("T3", "P120 D14 C24", BLUE_SOFT, BLUE),
    ("T4", "P120 D14 C12", AMBER_SOFT, AMBER),
    ("T5", "P120 D7  C12", GREEN_SOFT, GREEN),
    ("T6", "P120 D7  C12", ROSE_SOFT, ROSE),
    ("T7", "accept", MIST, NAVY),
]


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_backdrop(slide, primary: RGBColor, secondary: RGBColor) -> None:
    add_oval(slide, 10.00, -0.70, 4.10, 4.10, primary, primary, 0.5)
    add_oval(slide, -1.10, 5.55, 3.20, 3.20, secondary, secondary, 0.5)


def add_shape(slide, shape_type, x, y, w, h, fill_color, line_color=None, line_width=1.2):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color or fill_color
    shape.line.width = Pt(line_width)
    return shape


def add_round_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=1.2):
    return add_shape(slide, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h, fill_color, line_color, line_width)


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=1.2):
    return add_shape(slide, MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h, fill_color, line_color, line_width)


def add_oval(slide, x, y, w, h, fill_color, line_color=None, line_width=1.2):
    return add_shape(slide, MSO_AUTO_SHAPE_TYPE.OVAL, x, y, w, h, fill_color, line_color, line_width)


def add_line(slide, x1, y1, x2, y2, color, width=1.8):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=18,
    color=INK,
    bold=False,
    font_name=FONT_BODY,
    align=PP_ALIGN.LEFT,
    vertical=MSO_ANCHOR.TOP,
    margins=(0.06, 0.04, 0.04, 0.02),
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(margins[0])
    tf.margin_right = Inches(margins[1])
    tf.margin_top = Inches(margins[2])
    tf.margin_bottom = Inches(margins[3])
    tf.vertical_anchor = vertical

    for index, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = font_name
        p.font.color.rgb = color
    return shape


def add_card(slide, x, y, w, h, title, body, fill, accent, title_size=15.0, body_size=11.2):
    add_round_rect(slide, x, y, w, h, fill, accent, 1.5)
    add_textbox(slide, x + 0.16, y + 0.12, w - 0.32, 0.20, title, size=title_size, color=NAVY, bold=True)
    add_textbox(slide, x + 0.16, y + 0.40, w - 0.32, h - 0.50, body, size=body_size, color=INK)


def add_pill(slide, x, y, w, h, text, fill, text_color, line_color=None, font_size=10.0):
    add_round_rect(slide, x, y, w, h, fill, line_color or fill, 1)
    add_textbox(
        slide,
        x + 0.02,
        y + 0.05,
        w - 0.04,
        h - 0.08,
        text,
        size=font_size,
        color=text_color,
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
    )


def title_block(slide, chip, title, subtitle):
    add_pill(slide, 0.64, 0.44, 2.30, 0.30, chip, AMBER_SOFT, AMBER, AMBER, 9.6)
    add_textbox(slide, 0.64, 0.90, 8.20, 0.40, title, size=28, color=NAVY, bold=True, font_name=FONT_DISPLAY)
    add_textbox(slide, 0.66, 1.34, 9.20, 0.22, subtitle, size=12.0, color=MUTED)
    add_line(slide, 0.66, 1.68, 12.68, 1.68, SAND, 1.0)


def style_cell(cell, text, *, fill, color=INK, bold=False, size=11.0, align=PP_ALIGN.CENTER, font_name=FONT_BODY):
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.font.name = font_name
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color


def build_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, MIST, AMBER_SOFT)
    title_block(
        slide,
        "PACTUM CHALLENGE",
        "Negotiator",
        "Deterministic multi-issue negotiation across price, payment, delivery, and contract.",
    )

    add_textbox(
        slide,
        0.72,
        2.08,
        6.00,
        0.60,
        "The deck expands the math and the test evidence into separate slides so the decision system is readable end to end.",
        size=14.0,
        color=INK,
    )

    cards = [
        ("Commercial framing", "The engine models trade-offs across four deal terms instead of treating negotiation as price only.", BLUE_SOFT, BLUE),
        ("Technical framing", "AI can parse and draft language, but the authority to accept, counter, or reject stays rule-based.", TEAL_SOFT, TEAL),
        ("Evaluation framing", "The same formulas are then checked in StrategySimulationMatrixTest and PriceSensitivityTest.", GREEN_SOFT, GREEN),
    ]
    y = 3.02
    for title, body, fill, accent in cards:
        add_card(slide, 0.82, y, 5.86, 0.96, title, body, fill, accent)
        y += 1.14

    add_round_rect(slide, 7.24, 2.18, 5.02, 4.20, NAVY, NAVY, 1)
    add_textbox(slide, 7.56, 2.48, 2.60, 0.22, "Decision surface", size=19, color=WHITE, bold=True)
    add_oval(slide, 9.22, 3.52, 1.54, 1.54, WHITE, WHITE, 1)
    add_textbox(slide, 9.48, 3.96, 1.02, 0.48, "Buyer\nutility", size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    nodes = [
        ("Price", 8.04, 2.92, BLUE),
        ("Payment", 10.72, 2.92, TEAL),
        ("Delivery", 8.04, 5.08, GREEN),
        ("Contract", 10.72, 5.08, AMBER),
    ]
    for label, x, y, color in nodes:
        add_pill(slide, x, y, 1.14, 0.34, label, color, WHITE, color, 10.8)
    add_line(slide, 8.60, 3.24, 9.48, 3.90, WHITE, 2.0)
    add_line(slide, 11.28, 3.24, 10.50, 3.90, WHITE, 2.0)
    add_line(slide, 8.60, 5.08, 9.48, 4.96, WHITE, 2.0)
    add_line(slide, 11.28, 5.08, 10.50, 4.96, WHITE, 2.0)

    add_round_rect(slide, 7.58, 5.88, 4.24, 0.40, rgb("223242"), rgb("223242"), 1)
    add_textbox(slide, 7.76, 5.98, 3.88, 0.12, "Every round reduces to one comparable buyer-side score.", size=10.6, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, BLUE_SOFT, NAVY_SOFT)
    title_block(slide, "SYSTEM BOUNDARY", "Architecture", "Language can vary. Buyer commitments cannot.")

    add_pill(slide, 0.86, 2.10, 1.78, 0.28, "LLM assistance", BLUE_SOFT, BLUE, BLUE, 10.0)
    add_pill(slide, 7.42, 2.10, 2.02, 0.28, "Deterministic engine", GREEN_SOFT, GREEN, GREEN, 10.0)

    left_cards = [
        ("Parse supplier wording", "Free text becomes structured terms and supplier constraints."),
        ("Classify supplier intent", "Accept, select buyer option, propose new terms, reject, or unclear."),
        ("Draft buyer wording", "Natural language is generated only after the business decision is fixed."),
    ]
    y = 2.56
    for title, body in left_cards:
        add_card(slide, 0.82, y, 2.52, 1.08, title, body, BLUE_SOFT, BLUE, 13.8, 10.8)
        y += 1.32

    add_round_rect(slide, 4.28, 3.58, 1.86, 1.12, NAVY, NAVY, 1)
    add_textbox(slide, 4.48, 3.92, 1.46, 0.36, "Structured\noffer", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    right_cards = [
        ("Score offer", "Buyer utility across four issues."),
        ("Apply policy", "Round target, reservation checks, and settlement posture."),
        ("Choose action", "Accept, reject, or return one or more buyer-safe counters."),
    ]
    positions = [(7.30, 2.62), (9.94, 2.62), (8.62, 4.36)]
    for (title, body), (x, y) in zip(right_cards, positions):
        add_card(slide, x, y, 2.34, 1.10, title, body, GREEN_SOFT, GREEN, 14.2, 10.8)

    add_line(slide, 3.34, 3.10, 4.28, 4.12, NAVY, 2.0)
    add_line(slide, 3.34, 4.42, 4.28, 4.12, NAVY, 2.0)
    add_line(slide, 3.34, 5.74, 4.28, 4.12, NAVY, 2.0)
    add_line(slide, 6.14, 4.12, 7.30, 3.18, NAVY, 2.0)
    add_line(slide, 6.14, 4.12, 9.94, 3.18, NAVY, 2.0)
    add_line(slide, 8.48, 3.72, 9.24, 4.36, NAVY, 2.0)
    add_line(slide, 11.12, 3.72, 10.36, 4.36, NAVY, 2.0)

    add_round_rect(slide, 7.24, 5.94, 5.00, 0.54, WHITE, SAND, 1)
    add_textbox(slide, 7.42, 6.12, 4.64, 0.14, "AI parses and drafts. The engine scores, constrains, and decides.", size=11.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)


def build_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, NAVY_SOFT, TEAL_SOFT)
    title_block(slide, "CODE MAP", "Where The Logic Lives", "The math and the orchestration are split across small single-purpose classes.")

    table = slide.shapes.add_table(8, 2, Inches(0.84), Inches(2.16), Inches(7.20), Inches(4.64)).table
    table.columns[0].width = Inches(2.36)
    table.columns[1].width = Inches(4.84)
    headers = ["Class", "Role"]
    for col, header in enumerate(headers):
        style_cell(table.cell(0, col), header, fill=NAVY, color=WHITE, bold=True, size=11.0, align=PP_ALIGN.LEFT)

    rows = [
        ("NegotiationEngine", "Owns the records and enums: OfferVector, BuyerProfile, Context, Strategy, Response."),
        ("BuyerPreferenceScoring", "Normalizes each issue into a [0,1] buyer preference score."),
        ("BuyerUtilityCalculator", "Aggregates weighted issue scores into U_buyer."),
        ("DecisionMaker", "Computes round target utility and hard reject threshold."),
        ("StrategySettlementPolicy", "Adds strategy floors and price ceilings before acceptance."),
        ("CounterOfferGenerator", "Ranks gaps and generates next buyer-safe moves."),
        ("NegotiationApplicationService", "Handles session state, supplier intent, and explicit acceptance semantics."),
    ]
    for row_index, (name, role) in enumerate(rows, start=1):
        fill = WHITE if row_index % 2 else MIST
        style_cell(table.cell(row_index, 0), name, fill=fill, color=BLUE if row_index < 4 else TEAL, bold=True, size=10.4, align=PP_ALIGN.LEFT, font_name=FONT_MONO)
        style_cell(table.cell(row_index, 1), role, fill=fill, color=INK, size=10.2, align=PP_ALIGN.LEFT)

    records = [
        ("OfferVector", "price, paymentDays, deliveryDays, contractMonths", BLUE_SOFT, BLUE),
        ("IssueWeights", "relative issue importance before normalization", TEAL_SOFT, TEAL),
        ("BuyerProfile", "ideal offer + reservation offer + weights + reservationUtility", GREEN_SOFT, GREEN),
        ("NegotiationContext", "round, maxRounds, strategy, state, history", AMBER_SOFT, AMBER),
    ]
    y = 2.24
    for title, body, fill, accent in records:
        add_card(slide, 8.46, y, 4.02, 0.86, title, body, fill, accent, 14.0, 10.5)
        y += 1.02

    add_round_rect(slide, 8.46, 6.20, 4.02, 0.34, NAVY, NAVY, 1)
    add_textbox(slide, 8.62, 6.28, 3.70, 0.12, "The engine decision is a function of these records only.", size=10.0, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, TEAL_SOFT, MIST)
    title_block(slide, "MATHEMATICAL MODEL", "Issue Normalization", "Each issue is mapped onto the same buyer-side scale before weighting.")

    formulas = [
        ("Price", "priceScore = clamp((reservationPrice - offerPrice) / (reservationPrice - idealPrice))", BLUE_SOFT, BLUE),
        ("Payment", "paymentScore = clamp((offerPayment - reservationPayment) / (idealPayment - reservationPayment))", TEAL_SOFT, TEAL),
        ("Delivery", "deliveryScore = clamp((reservationDelivery - offerDelivery) / (reservationDelivery - idealDelivery))", GREEN_SOFT, GREEN),
        ("Contract", "contractScore = clamp((reservationContract - offerContract) / (reservationContract - idealContract))", AMBER_SOFT, AMBER),
    ]
    positions = [(0.84, 2.18), (6.70, 2.18), (0.84, 4.04), (6.70, 4.04)]
    for (title, formula, fill, accent), (x, y) in zip(formulas, positions):
        add_round_rect(slide, x, y, 5.76, 1.48, WHITE, SAND, 1.2)
        add_pill(slide, x + 0.16, y + 0.14, 1.18, 0.24, title, fill, accent, accent, 10.0)
        add_textbox(slide, x + 0.16, y + 0.54, 5.38, 0.34, formula, size=10.8, color=INK, font_name=FONT_MONO)
        add_textbox(slide, x + 0.18, y + 0.98, 5.34, 0.18, "Higher score always means better for the buyer.", size=9.6, color=MUTED)

    add_round_rect(slide, 0.84, 5.94, 11.62, 0.56, NAVY, NAVY, 1)
    add_textbox(slide, 1.06, 6.08, 11.18, 0.18, "Direction matters: lower price is better, longer payment is better, faster delivery is better, shorter contract is better.", size=11.2, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


def build_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, NAVY_SOFT, AMBER_SOFT)
    title_block(slide, "UTILITY FUNCTION", "Weighted Aggregation", "After normalization, the engine turns four issue scores into one buyer utility value.")

    add_round_rect(slide, 0.82, 2.10, 5.78, 3.96, WHITE, SAND, 1.2)
    add_textbox(slide, 1.08, 2.30, 2.40, 0.18, "Formula", size=17, color=NAVY, bold=True)
    add_textbox(
        slide,
        1.08,
        2.76,
        5.26,
        0.68,
        "U_buyer = w_price * priceScore\n"
        "        + w_payment * paymentScore\n"
        "        + w_delivery * deliveryScore\n"
        "        + w_contract * contractScore",
        size=14.0,
        color=INK,
        bold=True,
        font_name=FONT_MONO,
    )
    add_textbox(slide, 1.08, 3.70, 5.14, 0.18, "Weights are normalized first, so only their proportions matter.", size=10.2, color=MUTED)

    pieces = [
        ("Default weights", "price 0.40 | payment 0.20 | delivery 0.25 | contract 0.15", BLUE_SOFT, BLUE),
        ("Example offer", "price 108 | payment 45 | delivery 14 | contract 12", TEAL_SOFT, TEAL),
        ("Approx scores", "price 0.40 | payment 0.50 | delivery 0.70 | contract 0.67", GREEN_SOFT, GREEN),
    ]
    y = 4.18
    for title, body, fill, accent in pieces:
        add_card(slide, 1.04, y, 5.34, 0.54, title, body, fill, accent, 11.8, 9.8)
        y += 0.62

    add_round_rect(slide, 6.94, 2.10, 5.54, 3.96, WHITE, SAND, 1.2)
    add_textbox(slide, 7.20, 2.30, 2.80, 0.18, "Worked example", size=17, color=NAVY, bold=True)
    add_textbox(
        slide,
        7.20,
        2.74,
        5.00,
        1.08,
        "U = 0.40*0.40 + 0.20*0.50 + 0.25*0.70 + 0.15*0.67\n"
        "  = 0.160 + 0.100 + 0.175 + 0.100\n"
        "  = 0.535 approx",
        size=13.0,
        color=INK,
        font_name=FONT_MONO,
    )
    add_round_rect(slide, 7.20, 4.18, 4.96, 0.62, NAVY_SOFT, BLUE, 1.2)
    add_textbox(slide, 7.38, 4.34, 4.60, 0.22, "The supplier offer is not terrible, but it is still below a strong early-round target.", size=10.4, color=NAVY, bold=True)
    add_round_rect(slide, 7.20, 5.04, 4.96, 0.62, NAVY, NAVY, 1)
    add_textbox(slide, 7.40, 5.22, 4.56, 0.20, "This additive model is simple, explainable, and easy to tune.", size=10.6, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, BLUE_SOFT, NAVY_SOFT)
    title_block(slide, "DECISION LOGIC", "Acceptance Pipeline", "The engine evaluates offers in layers instead of one giant rule.")

    steps = [
        ("1. Reservation check", "If the offer is too far outside buyer limits, reject immediately.", BLUE_SOFT, BLUE),
        ("2. Buyer utility", "Calculate U_buyer from the normalized weighted model.", TEAL_SOFT, TEAL),
        ("3. Round target", "targetUtility = max(strategyCurve(round), reservationUtility).", AMBER_SOFT, AMBER),
        ("4. Settlement posture", "minimumUtility and maximumPrice add a stricter settlement gate.", GREEN_SOFT, GREEN),
        ("5. Decision", "Accept, counter, or reject based on those combined conditions.", ROSE_SOFT, ROSE),
    ]
    x = 0.84
    for index, (title, body, fill, accent) in enumerate(steps):
        add_card(slide, x, 2.42, 2.30, 2.16, title, body, fill, accent, 13.4, 10.4)
        if index < len(steps) - 1:
            add_line(slide, x + 2.30, 3.48, x + 2.50, 3.48, STONE, 2)
        x += 2.48

    add_round_rect(slide, 0.84, 5.12, 5.90, 1.12, WHITE, SAND, 1.2)
    add_textbox(slide, 1.08, 5.34, 2.70, 0.18, "Core decision rule", size=16, color=NAVY, bold=True)
    add_textbox(
        slide,
        1.08,
        5.72,
        5.40,
        0.34,
        "if U_buyer >= target and price <= maxPrice -> ACCEPT\n"
        "else if final round -> REJECT\n"
        "else -> COUNTER",
        size=11.4,
        color=INK,
        font_name=FONT_MONO,
    )

    add_round_rect(slide, 7.00, 5.12, 5.48, 1.12, NAVY, NAVY, 1)
    add_textbox(slide, 7.24, 5.34, 2.70, 0.18, "Important nuance", size=16, color=WHITE, bold=True)
    add_textbox(slide, 7.24, 5.70, 5.02, 0.36, "An offer can clear the round target and still fail the strategy settlement policy if the price ceiling or utility floor is not met.", size=10.6, color=WHITE)


def build_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, GREEN_SOFT, TEAL_SOFT)
    title_block(slide, "STRATEGY MATH", "Curves And Settlement Thresholds", "DecisionMaker and StrategySettlementPolicy turn business posture into numeric thresholds.")

    add_round_rect(slide, 0.82, 2.12, 6.06, 4.20, WHITE, SAND, 1.2)
    add_textbox(slide, 1.08, 2.32, 2.40, 0.18, "Target curves", size=17, color=NAVY, bold=True)
    target_rows = [
        ("Baseline", "T(p) = 1 - p", BLUE),
        ("Meso", "T(p) = 1 - p^1.35", AMBER),
        ("Boulware", "T(p) = 1 - p^2.4", GREEN),
        ("Conceder", "T(p) = 1 - sqrt(p)", ROSE),
        ("Tit for Tat", "T(p) = baseline - reciprocityBonus + firmnessPenalty", TEAL),
    ]
    y = 2.74
    for name, formula, color in target_rows:
        add_round_rect(slide, 1.02, y, 5.64, 0.48, WHITE, color, 1.2)
        add_textbox(slide, 1.20, y + 0.08, 1.38, 0.14, name, size=11.6, color=color, bold=True)
        add_textbox(slide, 2.62, y + 0.08, 3.80, 0.14, formula, size=10.8, color=INK, font_name=FONT_MONO)
        y += 0.58

    add_round_rect(slide, 1.02, 5.74, 5.64, 0.38, NAVY, NAVY, 1)
    add_textbox(slide, 1.20, 5.84, 5.26, 0.14, "Hard reject thresholds: Baseline 0.1200 | Meso 0.1000 | Boulware 0.1500 | Conceder 0.0800", size=9.8, color=WHITE, align=PP_ALIGN.CENTER)

    add_round_rect(slide, 7.08, 2.12, 5.42, 4.20, WHITE, SAND, 1.2)
    add_textbox(slide, 7.34, 2.32, 2.90, 0.18, "Settlement policy", size=17, color=NAVY, bold=True)
    add_textbox(slide, 7.34, 2.72, 4.92, 0.34, "minimumUtility = max(reservationUtility, strategyFloor)\nmaximumPrice  = idealPrice + priceSpan * priceRatio", size=11.0, color=INK, font_name=FONT_MONO)

    thresholds = [
        ("Boulware", "priceRatio 0.55 | utilityFloor 0.7100", GREEN),
        ("Baseline", "priceRatio 0.70 | utilityFloor 0.6200", BLUE),
        ("Tit for Tat", "priceRatio 0.62 + bonus | utilityFloor 0.6200 - relief", TEAL),
        ("Meso", "priceRatio 0.78 | utilityFloor 0.6000", AMBER),
        ("Conceder", "priceRatio 1.00 | utilityFloor 0.5000", ROSE),
    ]
    y = 3.36
    for name, body, color in thresholds:
        add_round_rect(slide, 7.28, y, 5.04, 0.50, WHITE, color, 1.2)
        add_textbox(slide, 7.46, y + 0.08, 1.30, 0.14, name, size=11.4, color=color, bold=True)
        add_textbox(slide, 8.82, y + 0.08, 3.30, 0.14, body, size=10.2, color=INK, font_name=FONT_MONO)
        y += 0.60

    add_round_rect(slide, 7.28, 5.74, 5.04, 0.38, NAVY, NAVY, 1)
    add_textbox(slide, 7.42, 5.84, 4.76, 0.14, "The same utility can settle under Conceder and still counter under Boulware.", size=9.8, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, AMBER_SOFT, MIST)
    title_block(slide, "COUNTEROFFER ENGINE", "Gap Ranking And MESO", "CounterOfferGenerator selects the next move by weighted gap, not by random variation.")

    add_round_rect(slide, 0.82, 2.14, 5.82, 4.12, WHITE, SAND, 1.2)
    add_textbox(slide, 1.08, 2.34, 2.40, 0.18, "Gap ranking", size=17, color=NAVY, bold=True)
    add_textbox(
        slide,
        1.08,
        2.74,
        5.28,
        0.38,
        "gap_i = (distance from supplier offer to buyer ideal / issue span) * normalizedWeight_i",
        size=11.4,
        color=INK,
        font_name=FONT_MONO,
    )
    steps = [
        ("1. Rank issues", "Find the most valuable remaining gap for the buyer.", BLUE_SOFT, BLUE),
        ("2. Respect constraints", "Blocked issues are zeroed if supplier floors or ceilings stop movement.", TEAL_SOFT, TEAL),
        ("3. Move halfway", "The chosen issue moves toward the buyer ideal with a bounded midpoint step.", GREEN_SOFT, GREEN),
        ("4. Tune result", "Rebalance price, keep history consistency, and clamp to settlement policy.", AMBER_SOFT, AMBER),
    ]
    y = 3.34
    for title, body, fill, accent in steps:
        add_card(slide, 1.02, y, 5.40, 0.60, title, body, fill, accent, 11.8, 9.8)
        y += 0.72

    add_round_rect(slide, 7.02, 2.14, 5.48, 4.12, WHITE, SAND, 1.2)
    add_textbox(slide, 7.28, 2.34, 2.10, 0.18, "MESO behavior", size=17, color=NAVY, bold=True)
    add_textbox(slide, 7.28, 2.70, 4.90, 0.20, "MESO keeps the same safety checks but can return up to three viable options in one round.", size=10.6, color=INK)

    options = [
        ("Option 1", "price-focused move", BLUE_SOFT, BLUE),
        ("Option 2", "delivery-focused move", GREEN_SOFT, GREEN),
        ("Option 3", "contract-focused move", AMBER_SOFT, AMBER),
    ]
    y = 3.20
    for name, body, fill, accent in options:
        add_round_rect(slide, 7.30, y, 4.90, 0.56, fill, accent, 1.2)
        add_textbox(slide, 7.48, y + 0.10, 0.90, 0.14, name, size=11.0, color=accent, bold=True)
        add_textbox(slide, 8.44, y + 0.10, 3.52, 0.14, body, size=10.0, color=INK)
        y += 0.72

    add_round_rect(slide, 7.30, 5.36, 4.90, 0.60, NAVY, NAVY, 1)
    add_textbox(slide, 7.52, 5.54, 4.46, 0.18, "The supplier can then reveal preference without the buyer weakening its position.", size=10.4, color=WHITE, align=PP_ALIGN.CENTER)


def build_slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, GREEN_SOFT, BLUE_SOFT)
    title_block(slide, "TEST EVIDENCE", "StrategySimulationMatrixTest", "Seven deterministic scenarios confirm that the strategies remain behaviorally distinct.")

    table = slide.shapes.add_table(8, 6, Inches(0.80), Inches(2.14), Inches(8.02), Inches(4.08)).table
    table.columns[0].width = Inches(2.44)
    for col in range(1, 6):
        table.columns[col].width = Inches(1.116)
    headers = ["Scenario", "Baseline", "Meso", "Boulware", "Conceder", "Tit for Tat"]
    for col, header in enumerate(headers):
        style_cell(table.cell(0, col), header, fill=NAVY, color=WHITE, bold=True, size=10.4, align=PP_ALIGN.LEFT if col == 0 else PP_ALIGN.CENTER)
    for row, (scenario, outcomes) in enumerate(SCENARIO_MATRIX, start=1):
        fill = WHITE if row % 2 else MIST
        style_cell(table.cell(row, 0), scenario, fill=fill, color=INK, bold=True, size=10.0, align=PP_ALIGN.LEFT)
        for col, outcome in enumerate(outcomes, start=1):
            accepted = outcome.startswith("A")
            style_cell(table.cell(row, col), outcome, fill=GREEN_SOFT if accepted else ROSE_SOFT, color=GREEN if accepted else ROSE, bold=True, size=12.0)

    add_pill(slide, 0.82, 6.34, 1.24, 0.24, "A = accepted", GREEN_SOFT, GREEN, GREEN, 9.2)
    add_pill(slide, 2.18, 6.34, 1.18, 0.24, "R = rejected", ROSE_SOFT, ROSE, ROSE, 9.2)
    add_textbox(slide, 3.52, 6.36, 5.10, 0.14, "number = round where the scenario closed or failed", size=9.4, color=MUTED)

    add_round_rect(slide, 9.08, 2.14, 3.92, 2.26, WHITE, SAND, 1.2)
    add_textbox(slide, 9.34, 2.34, 2.60, 0.18, "Accepted averages", size=16, color=NAVY, bold=True)
    avg_table = slide.shapes.add_table(6, 3, Inches(9.24), Inches(2.68), Inches(3.56), Inches(1.52)).table
    avg_table.columns[0].width = Inches(1.68)
    avg_table.columns[1].width = Inches(0.90)
    avg_table.columns[2].width = Inches(0.90)
    for col, header in enumerate(["Strategy", "Buyer U", "Avg R"]):
        style_cell(avg_table.cell(0, col), header, fill=NAVY, color=WHITE, bold=True, size=10.0)
    for row, (name, utility, round_value, color) in enumerate(STRATEGY_ROWS, start=1):
        fill = WHITE if row % 2 else MIST
        style_cell(avg_table.cell(row, 0), name, fill=fill, color=color, bold=True, size=9.8, align=PP_ALIGN.LEFT)
        style_cell(avg_table.cell(row, 1), utility, fill=fill, color=INK, size=9.8)
        style_cell(avg_table.cell(row, 2), round_value, fill=fill, color=INK, size=9.8)

    add_round_rect(slide, 9.08, 4.56, 3.92, 1.66, WHITE, SAND, 1.2)
    add_textbox(slide, 9.30, 4.76, 2.20, 0.18, "What the test proves", size=15, color=NAVY, bold=True)
    bullets = [
        ("Boulware preserves the most buyer utility.", GREEN),
        ("Baseline stays in the middle.", BLUE),
        ("Conceder settles earliest.", ROSE),
        ("Meso differs through exploration and later closure.", AMBER),
    ]
    y = 5.12
    for text, color in bullets:
        add_oval(slide, 9.28, y + 0.03, 0.12, 0.12, color, color, 1)
        add_textbox(slide, 9.48, y, 3.20, 0.14, text, size=9.8, color=INK)
        y += 0.24


def build_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, PAPER)
    add_backdrop(slide, BLUE_SOFT, ROSE_SOFT)
    title_block(slide, "TEST EVIDENCE", "PriceSensitivityTest", "The supplier script is fixed, so the closing price differences come entirely from buyer strategy.")

    add_round_rect(slide, 0.82, 2.12, 12.00, 1.12, WHITE, SAND, 1.2)
    add_textbox(slide, 1.06, 2.30, 2.10, 0.18, "Fixed supplier flow", size=16, color=NAVY, bold=True)
    add_pill(slide, 3.24, 2.26, 2.16, 0.24, "Price 120 fixed", BLUE_SOFT, BLUE, BLUE, 9.8)
    add_pill(slide, 5.50, 2.26, 2.22, 0.24, "Payment 60 fixed", TEAL_SOFT, TEAL, TEAL, 9.8)
    add_pill(slide, 7.86, 2.26, 3.04, 0.24, "Delivery and contract improve", AMBER_SOFT, AMBER, AMBER, 9.8)

    xs = [0.96, 2.62, 4.28, 5.94, 7.60, 9.26, 10.92]
    for left, right in zip(xs, xs[1:]):
        add_line(slide, left + 1.38, 2.96, right - 0.08, 2.96, STONE, 1.6)
    for (turn, label, fill, accent), x in zip(SCRIPT_STEPS, xs):
        add_round_rect(slide, x, 2.72, 1.46, 0.42, fill, accent, 1.2)
        add_textbox(slide, x + 0.08, 2.80, 0.30, 0.12, turn, size=9.0, color=accent, bold=True)
        add_textbox(slide, x + 0.40, 2.80, 0.96, 0.12, label, size=8.8, color=INK, font_name=FONT_MONO)

    add_round_rect(slide, 0.82, 3.54, 7.40, 3.02, WHITE, SAND, 1.2)
    add_textbox(slide, 1.06, 3.74, 3.20, 0.18, "Average closing price", size=16, color=NAVY, bold=True)
    add_textbox(slide, 6.02, 3.76, 1.30, 0.12, "close | first", size=9.4, color=MUTED, align=PP_ALIGN.CENTER)

    axis_left = 2.48
    axis_right = 6.22
    axis_y = 6.18
    add_line(slide, axis_left, axis_y, axis_right, axis_y, STONE, 1.6)
    for value in [100, 105, 110, 115, 120]:
        x = axis_left + (value - 100) / 20.0 * (axis_right - axis_left)
        add_line(slide, x, 4.10, x, axis_y, GRID, 0.8)
        add_textbox(slide, x - 0.16, axis_y + 0.04, 0.32, 0.12, str(value), size=9.0, color=MUTED, align=PP_ALIGN.CENTER)

    y = 4.22
    for label, value, close_round, first_offer_round, color in PRICE_RESULTS:
        add_oval(slide, 1.06, y + 0.03, 0.12, 0.12, color, color, 1)
        add_textbox(slide, 1.24, y, 1.16, 0.14, label, size=10.0, color=color, bold=True)
        width = max(0.26, (value - 100.0) / 20.0 * (axis_right - axis_left))
        add_round_rect(slide, axis_left, y, width, 0.22, color, color, 1)
        add_textbox(slide, axis_left + width + 0.06, y + 0.02, 0.56, 0.12, f"{value:.2f}", size=9.8, color=INK)
        add_textbox(slide, 5.76, y + 0.02, 1.46, 0.12, f"{close_round} | {first_offer_round}", size=9.2, color=MUTED, align=PP_ALIGN.CENTER)
        y += 0.40

    add_textbox(slide, 1.06, 6.34, 6.90, 0.14, "Current passing output: Boulware 106.50, Baseline 111.00, Tit for Tat 113.33, Meso 113.40, Conceder 120.00.", size=9.6, color=MUTED)

    add_round_rect(slide, 8.44, 3.54, 4.38, 3.02, NAVY, NAVY, 1)
    add_textbox(slide, 8.70, 3.74, 2.20, 0.18, "Interpretation", size=16, color=WHITE, bold=True)
    notes = [
        ("Boulware", "holds price hardest and still closes", GREEN),
        ("Baseline", "gives some price to keep momentum", BLUE),
        ("Tit for Tat", "tracks close to Baseline", TEAL),
        ("Meso", "primary offer closes at 113.40", AMBER),
        ("Conceder", "accepts the ceiling to close", ROSE),
    ]
    y = 4.18
    for label, body, color in notes:
        add_round_rect(slide, 8.64, y, 3.98, 0.42, WHITE, WHITE, 1)
        add_oval(slide, 8.78, y + 0.14, 0.12, 0.12, color, color, 1)
        add_textbox(slide, 9.00, y + 0.10, 1.12, 0.12, label, size=10.0, color=color, bold=True)
        add_textbox(slide, 10.12, y + 0.10, 2.30, 0.12, body, size=9.4, color=INK)
        y += 0.50

    add_round_rect(slide, 8.64, 6.16, 3.98, 0.26, rgb("223242"), rgb("223242"), 1)
    add_textbox(slide, 8.82, 6.22, 3.62, 0.10, "Source: PriceSensitivityTest over 1,000 runs", size=9.0, color=WHITE, align=PP_ALIGN.CENTER)


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)
    build_slide_7(prs)
    build_slide_8(prs)
    build_slide_9(prs)
    build_slide_10(prs)

    OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PPTX)
    print(f"Saved to {OUT_PPTX}")


if __name__ == "__main__":
    build_presentation()
